from pptx import Presentation


def get_tables(prs):
    idx_tables = {}
    slide_idx = 0

    for slide in prs.slides:            
        if len(idx_tables) == 9:
            break            
        for shape in slide.shapes:
            if not shape.has_table:
                continue        
            for cell in shape.table.iter_cells():            
                if 'TCV' in cell.text:
                    if 'TCV' not in idx_tables:                            
                        idx_tables['TCV'] = {}
                        idx_tables['TCV']['slide_idx'] = slide_idx
                        idx_tables['TCV']['shape_id'] = shape.shape_id
                        break
                if 'THC' in cell.text:
                    if 'THC' not in idx_tables:
                        idx_tables['THC'] = {}
                        idx_tables['THC']['slide_idx'] = slide_idx
                        idx_tables['THC']['shape_id'] = shape.shape_id
                        break
                if 'TQF' in cell.text:
                    if 'TQF' not in idx_tables:
                        idx_tables['TQF'] = {}
                        idx_tables['TQF']['slide_idx'] = slide_idx
                        idx_tables['TQF']['shape_id'] = shape.shape_id
                        break
                if 'ARMAS' in cell.text:
                    if 'ARMAS' not in idx_tables:
                        idx_tables['IAF'] = {}
                        idx_tables['IAF']['slide_idx'] = slide_idx
                        idx_tables['IAF']['shape_id'] = shape.shape_id
                        break
                if 'AUTORES DE CV PRESOS' in cell.text:
                    if 'TRI' not in idx_tables:
                        idx_tables['TRI'] = {}
                        idx_tables['TRI']['slide_idx'] = slide_idx
                        idx_tables['TRI']['shape_id'] = shape.shape_id
                        break
                if 'EFETIVIDADE' in cell.text:
                    if 'EFETIVIDADE' not in idx_tables:
                        idx_tables['RQV'] = {}
                        idx_tables['RQV']['slide_idx'] = slide_idx
                        idx_tables['RQV']['shape_id'] = shape.shape_id
                        break
                if 'DENÚNCIAS' in cell.text:
                    if 'DENÚNCIAS' not in idx_tables:
                        idx_tables['DDU'] = {}
                        idx_tables['DDU']['slide_idx'] = slide_idx
                        idx_tables['DDU']['shape_id'] = shape.shape_id
                        break
                if 'TOLS' in cell.text:
                    if 'TOLS' not in idx_tables:
                        idx_tables['OLS'] = {}
                        idx_tables['OLS']['slide_idx'] = slide_idx
                        idx_tables['OLS']['shape_id'] = shape.shape_id
                        break
                if 'QTD OP' in cell.text:
                    if 'QTD OP' not in idx_tables:
                        idx_tables['POG'] = {}
                        idx_tables['POG']['slide_idx'] = slide_idx
                        idx_tables['POG']['shape_id'] = shape.shape_id
                        break
        slide_idx+=1
    
    tables = {}
    for k, v in idx_tables.items():
        slide = prs.slides[v['slide_idx']]
        shapes = slide.shapes
        for shape in shapes:
            if shape.shape_id == v['shape_id']:
                tables[k] = shape
    
    return tables


def replace_text(MESES, mes, ano):
    prs = Presentation('files/modelo_apresentacao/Modelo (Novo) GDO 23º BPM.pptx')    
    for shape in prs.slides[0].shapes:
        if 'Mês' in shape.text:
            for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            p = paragraph._p  # the lxml element containing the `<a:p>` paragraph element
        #                            remove all but the first run
                            for idx, run in enumerate(paragraph.runs):
                                if idx == 0:
                                    continue
                                p.remove(run._r)
                            cur_text = run.text
                            new_text = cur_text.replace('Mês de Ano', f'{MESES[mes]} de {ano}')
                            run.text = new_text
    prs.save(f'files/{ano}/{mes}/Apresentação GDO - {MESES[mes]} de {ano}.pptx')


def replace_text_table_header(path_name_source_file, path_name_result_file, replace_dict):
    prs = Presentation(path_name_source_file)
    for slide in prs.slides:
        for shape in slide.shapes:
                for match, replacement in replace_dict.items():
                    if shape.has_table:
                        print(dir(shape.table.rows))
#                         for cell in shape.table.iter_cells():
#                         for cell in shape.table.rows:
#                             print(cell.text)
#                             if match == cell.text:
#                                 for paragraph in cell.text_frame.paragraphs:
#                                    for run in paragraph.runs:
#                                        p = paragraph._p  # the lxml element containing the `<a:p>` paragraph element
#         #                            remove all but the first run
#                                        for idx, run in enumerate(paragraph.runs):
#                                            if idx == 0:
#                                                continue
#                                            p.remove(run._r)
#                                        cur_text = run.text
#                                        new_text = cur_text.replace(str(match), str(replacement))
#                                        run.text = new_text
#     prs.save(path_name_result_file)

